# ========= Copyright 2023-2024 @ CAMEL-AI.org. All Rights Reserved. =========
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
# ========= Copyright 2023-2024 @ CAMEL-AI.org. All Rights Reserved. =========


import glob
import os
import random
import re
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING, Any, Dict, List, Optional, Tuple, Union

import numpy as np
import pptx
import requests
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from sentence_transformers import SentenceTransformer

if TYPE_CHECKING:
    from pptx import presentation
    from pptx.slide import Slide
    from pptx.text.text import TextFrame

from camel.logger import get_logger
from camel.toolkits.base import BaseToolkit
from camel.toolkits.function_tool import FunctionTool
from camel.utils import MCPServer, api_keys_required

logger = get_logger(__name__)

# Constants
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400

STEP_BY_STEP_PROCESS_MARKER = '>> '
ICON_BEGINNING_MARKER = '[['
ICON_END_MARKER = ']]'

ICON_SIZE = Inches(0.8)
ICON_BG_SIZE = Inches(1)

SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")
BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')


@MCPServer()
class PPTXToolkit(BaseToolkit):
    r"""A toolkit for creating and writing PowerPoint presentations (PPTX
    files).

    This class provides cross-platform support for creating PPTX files with
    title slides, content slides, text formatting, and image embedding.
    """

    def __init__(
        self,
        output_dir: str = "./",
        timeout: Optional[float] = None,
    ) -> None:
        r"""Initialize the PPTXToolkit.

        Args:
            output_dir (str): The default directory for output files.
                Defaults to the current working directory.
            timeout (Optional[float]): The timeout for the toolkit.
                (default: :obj: `None`)
        """
        super().__init__(timeout=timeout)
        self.output_dir = Path(output_dir).resolve()
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(
            f"PPTXToolkit initialized with output directory: {self.output_dir}"
        )
        # For semantic search over local PNG icons
        self._png_icon_dir = os.path.abspath(
            os.path.join(os.path.dirname(__file__), './png128')
        )
        self._png_icon_index = None  # List of (icon_name, full_path)
        self._png_icon_embeddings = None
        self._png_icon_model = None

    def _build_png_icon_index(self):
        r"""Build index and embeddings for all PNG icons in the icon
        directory."""
        if (
            self._png_icon_index is not None
            and self._png_icon_embeddings is not None
        ):
            return
        icon_files = glob.glob(os.path.join(self._png_icon_dir, '*.png'))
        icon_names = [
            os.path.splitext(os.path.basename(f))[0] for f in icon_files
        ]
        self._png_icon_index = list(zip(icon_names, icon_files))
        # Use a sentence transformer for semantic search
        if self._png_icon_model is None:
            self._png_icon_model = SentenceTransformer('all-MiniLM-L6-v2')
        self._png_icon_embeddings = self._png_icon_model.encode(icon_names)

    def _find_closest_png_icon(
        self, query: str, threshold: float = 0.5
    ) -> Optional[str]:
        r"""Find the closest matching PNG icon using semantic search.
        Returns the file path or None."""
        self._build_png_icon_index()
        if not self._png_icon_index or not self._png_icon_embeddings:
            return None
        # Ensure _png_icon_embeddings is a numpy array for mypy
        png_icon_embeddings = self._png_icon_embeddings  # type: ignore[assignment]
        icon_names, icon_files = zip(*self._png_icon_index)
        query_emb = self._png_icon_model.encode([query])[0]
        similarities = np.dot(png_icon_embeddings, query_emb) / (
            np.linalg.norm(png_icon_embeddings, axis=1)
            * np.linalg.norm(query_emb)
        )
        best_idx = int(np.argmax(similarities))
        if similarities[best_idx] >= threshold:
            return icon_files[best_idx]
        return None

    def _resolve_filepath(self, file_path: str) -> Path:
        r"""Convert the given string path to a Path object.

        If the provided path is not absolute, it is made relative to the
        default output directory. The filename part is sanitized to replace
        spaces and special characters with underscores, ensuring safe usage
        in downstream processing.

        Args:
            file_path (str): The file path to resolve.

        Returns:
            Path: A fully resolved (absolute) and sanitized Path object.
        """
        path_obj = Path(file_path)
        if not path_obj.is_absolute():
            path_obj = self.output_dir / path_obj

        sanitized_filename = self._sanitize_filename(path_obj.name)
        path_obj = path_obj.parent / sanitized_filename
        return path_obj.resolve()

    def _sanitize_filename(self, filename: str) -> str:
        r"""Sanitize a filename by replacing special characters and spaces.

        Args:
            filename (str): The filename to sanitize.

        Returns:
            str: The sanitized filename.
        """
        import re

        # Replace spaces and special characters with underscores
        sanitized = re.sub(r'[^\w\-_\.]', '_', filename)
        # Remove multiple consecutive underscores
        sanitized = re.sub(r'_+', '_', sanitized)
        return sanitized

    def _format_text(
        self, frame_paragraph, text: str, set_color_to_white=False
    ) -> None:
        r"""Apply bold and italic formatting while preserving the original
        word order.

        Args:
            frame_paragraph: The paragraph to format.
            text (str): The text to format.
            set_color_to_white (bool): Whether to set the color to white.
                (default: :obj: `False`)
        """
        from pptx.dml.color import RGBColor

        matches = list(BOLD_ITALICS_PATTERN.finditer(text))
        last_index = 0

        for match in matches:
            start, end = match.span()
            if start > last_index:
                run = frame_paragraph.add_run()
                run.text = text[last_index:start]
                if set_color_to_white:
                    run.font.color.rgb = RGBColor(255, 255, 255)

            if match.group(2):  # Bold
                run = frame_paragraph.add_run()
                run.text = match.group(2)
                run.font.bold = True
                if set_color_to_white:
                    run.font.color.rgb = RGBColor(255, 255, 255)
            elif match.group(3):  # Italics
                run = frame_paragraph.add_run()
                run.text = match.group(3)
                run.font.italic = True
                if set_color_to_white:
                    run.font.color.rgb = RGBColor(255, 255, 255)

            last_index = end

        if last_index < len(text):
            run = frame_paragraph.add_run()
            run.text = text[last_index:]
            if set_color_to_white:
                run.font.color.rgb = RGBColor(255, 255, 255)

    def _add_bulleted_items(
        self,
        text_frame: "TextFrame",
        flat_items_list: List[Tuple[str, int]],
        set_color_to_white: bool = False,
    ) -> None:
        r"""Add a list of texts as bullet points and apply formatting.

        Args:
            text_frame (TextFrame): The text frame where text is to be
                displayed.
            flat_items_list (List[Tuple[str, int]]): The list of items to be
                displayed.
            set_color_to_white (bool): Whether to set the font color to white.
                (default: :obj: `False`)
        """
        if not flat_items_list:
            logger.warning("Empty bullet point list provided")
            return
        for idx, item_content in enumerate(flat_items_list):
            item_text, item_level = item_content

            if idx == 0:
                if not text_frame.paragraphs:
                    # Ensure a paragraph exists if the frame is empty or
                    # cleared
                    paragraph = text_frame.add_paragraph()
                else:
                    # Use the first existing paragraph
                    paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()

            paragraph.level = item_level

            self._format_text(
                paragraph,
                item_text.removeprefix(STEP_BY_STEP_PROCESS_MARKER),
                set_color_to_white=set_color_to_white,
            )

    def _get_flat_list_of_contents(
        self, items: List[Union[str, List[Any]]], level: int
    ) -> List[Tuple[str, int]]:
        r"""Flatten a hierarchical list of bullet points to a single list.

        Args:
            items (List[Union[str, List[Any]]]): A bullet point (string or
                list).
            level (int): The current level of hierarchy.

        Returns:
            List[Tuple[str, int]]: A list of (bullet item text, hierarchical
                level) tuples.
        """
        flat_list = []

        for item in items:
            if isinstance(item, str):
                flat_list.append((item, level))
            elif isinstance(item, list):
                flat_list.extend(
                    self._get_flat_list_of_contents(item, level + 1)
                )

        return flat_list

    def _get_slide_width_height_inches(
        self, presentation: "presentation.Presentation"
    ) -> Tuple[float, float]:
        r"""Get the dimensions of a slide in inches.

        Args:
            presentation (presentation.Presentation): The presentation object.

        Returns:
            Tuple[float, float]: The width and height in inches.
        """
        slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * (
            presentation.slide_width or 0
        )
        slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * (
            presentation.slide_height or 0
        )
        return slide_width_inch, slide_height_inch

    def _write_pptx_file(
        self,
        file_path: Path,
        content: List[Dict[str, Any]],
        template: Optional[str] = None,
    ) -> None:
        r"""Write text content to a PPTX file with enhanced formatting.

        Args:
            file_path (Path): The target file path.
            content (List[Dict[str, Any]]): The content to write to the PPTX
                file. Must be a list of dictionaries where:
                - First element: Title slide with keys 'title' and 'subtitle'
                - Subsequent elements: Content slides with keys 'title', 'text'
            template (Optional[str]): The name of the template to use. If not
                provided, the default template will be used. (default: :obj:
                `None`)
        """
        import pptx

        # Use template if provided, otherwise create new presentation
        if template is not None:
            template_path = Path(template).resolve()
            if not template_path.exists():
                logger.warning(
                    f"Template file not found: {template_path}, using "
                    "default template"
                )
                presentation = pptx.Presentation()
            else:
                presentation = pptx.Presentation(str(template_path))
                # Clear all existing slides by removing them from the slide
                # list
                while len(presentation.slides) > 0:
                    rId = presentation.slides._sldIdLst[-1].rId
                    presentation.part.drop_rel(rId)
                    del presentation.slides._sldIdLst[-1]
        else:
            presentation = pptx.Presentation()

        slide_width_inch, slide_height_inch = (
            self._get_slide_width_height_inches(presentation)
        )

        # Process slides
        if content:
            # Title slide (first element)
            title_slide_data = content.pop(0) if content else {}
            title_layout = presentation.slide_layouts[0]
            title_slide = presentation.slides.add_slide(title_layout)

            # Set title and subtitle
            if title_slide.shapes.title:
                title_slide.shapes.title.text_frame.clear()
                self._format_text(
                    title_slide.shapes.title.text_frame.paragraphs[0],
                    title_slide_data.get("title", ""),
                )

            if len(title_slide.placeholders) > 1:
                subtitle = title_slide.placeholders[1]
                subtitle.text_frame.clear()
                self._format_text(
                    subtitle.text_frame.paragraphs[0],
                    title_slide_data.get("subtitle", ""),
                )

            # Content slides
            for slide_data in content:
                if not isinstance(slide_data, dict):
                    continue

                # Handle different slide types
                if 'table' in slide_data:
                    self._handle_table(
                        presentation,
                        slide_data,
                    )
                elif 'bullet_points' in slide_data:
                    if any(
                        step.startswith(ICON_BEGINNING_MARKER)
                        for step in slide_data['bullet_points']
                    ):
                        self._handle_icons_ideas(
                            presentation,
                            slide_data,
                            slide_width_inch,
                            slide_height_inch,
                        )
                    elif any(
                        step.startswith(STEP_BY_STEP_PROCESS_MARKER)
                        for step in slide_data['bullet_points']
                    ):
                        self._handle_step_by_step_process(
                            presentation,
                            slide_data,
                            slide_width_inch,
                            slide_height_inch,
                        )
                    else:
                        self._handle_default_display(
                            presentation,
                            slide_data,
                        )

        # Save the presentation
        presentation.save(str(file_path))
        logger.debug(f"Wrote PPTX to {file_path} with enhanced formatting")

    def create_presentation(
        self,
        content: str,
        filename: str,
        template: Optional[str] = None,
    ) -> str:
        r"""Create a PowerPoint presentation (PPTX) file.

        Args:
            content (str): The content to write to the PPTX file as a JSON
                string. Must represent a list of dictionaries with the
                following structure:
                - First dict: title slide {"title": str, "subtitle": str}
                - Other dicts: content slides, which can be one of:
                    * Bullet/step slides: {"heading": str, "bullet_points":
                    list of str or nested lists, "img_keywords": str
                    (optional)}
                        - If any bullet point starts with '>> ', it will be
                        rendered as a step-by-step process.
                        - "img_keywords" can be a URL or search keywords for
                        an image (optional).
                    * Table slides: {"heading": str, "table": {"headers": list
                    of str, "rows": list of list of str}}
            filename (str): The name or path of the file. If a relative path is
                supplied, it is resolved to self.output_dir.
            template (Optional[str]): The path to the template PPTX file.
                Initializes a presentation from a given template file Or PPTX
                file. (default: :obj: `None`)

        Returns:
            str: A success message indicating the file was created.

        Example:
        [
            {
                "title": "Presentation Title",
                "subtitle": "Presentation Subtitle"
            },
            {
                "heading": "Slide Title",
                "bullet_points": [
                    "**Bold text** for emphasis",
                    "*Italic text* for additional emphasis",
                    "Regular text for normal content"
                ],
                "img_keywords": "relevant search terms for images"
            },
            {
                "heading": "Step-by-Step Process",
                "bullet_points": [
                    ">> **Step 1:** First step description",
                    ">> **Step 2:** Second step description",
                    ">> **Step 3:** Third step description"
                ],
                "img_keywords": "process workflow steps"
            },
            {
                "heading": "Comparison Table",
                "table": {
                    "headers": ["Column 1", "Column 2", "Column 3"],
                    "rows": [
                        ["Row 1, Col 1", "Row 1, Col 2", "Row 1, Col 3"],
                        ["Row 2, Col 1", "Row 2, Col 2", "Row 2, Col 3"]
                    ]
                },
                "img_keywords": "comparison visualization"
            }
        ]
        """
        # Ensure filename has .pptx extension
        if not filename.lower().endswith('.pptx'):
            filename += '.pptx'

        # Resolve file path
        file_path = self._resolve_filepath(filename)

        # Parse and validate content format
        try:
            import json

            parsed_content = json.loads(content)
        except json.JSONDecodeError as e:
            logger.error(f"Content must be valid JSON: {e}")
            return "Failed to parse content as JSON"

        if not isinstance(parsed_content, list):
            logger.error(
                f"PPTX content must be a list of dictionaries, "
                f"got {type(parsed_content).__name__}"
            )
            return "PPTX content must be a list of dictionaries"

        try:
            # Create the PPTX file
            self._write_pptx_file(file_path, parsed_content.copy(), template)

            success_msg = (
                f"PowerPoint presentation successfully created: {file_path}"
            )
            logger.info(success_msg)
            return success_msg

        except Exception as e:
            error_msg = f"Failed to create PPTX file {file_path}: {e!s}"
            logger.error(error_msg)
            return error_msg

    def _handle_default_display(
        self,
        presentation: "presentation.Presentation",
        slide_json: Dict[str, Any],
    ) -> None:
        r"""Display a list of text in a slide.

        Args:
            presentation (presentation.Presentation): The presentation object.
            slide_json (Dict[str, Any]): The content of the slide as JSON data.
        """
        if 'img_keywords' in slide_json:
            if self._handle_display_image__in_foreground(
                presentation,
                slide_json,
            ):
                return

        # Image display failed, so display only text
        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes
        title_shape = shapes.title

        try:
            body_shape = shapes.placeholders[1]
        except KeyError:
            # Get placeholders from the slide without layout_number
            placeholders = self._get_slide_placeholders(slide)
            body_shape = shapes.placeholders[placeholders[0][0]]

        title_shape.text = self._remove_slide_number_from_heading(
            slide_json['heading']
        )
        text_frame = body_shape.text_frame

        flat_items_list = self._get_flat_list_of_contents(
            slide_json['bullet_points'], level=0
        )
        self._add_bulleted_items(text_frame, flat_items_list)

    @api_keys_required(
        [
            ("api_key", 'PEXELS_API_KEY'),
        ]
    )
    def _handle_display_image__in_foreground(
        self,
        presentation: "presentation.Presentation",
        slide_json: Dict[str, Any],
    ) -> bool:
        r"""Create a slide with text and image using a picture placeholder
        layout.

        Args:
            presentation (presentation.Presentation): The presentation object.
            slide_json (Dict[str, Any]): The content of the slide as JSON data.

        Returns:
            bool: True if the slide has been processed.
        """

        img_keywords = slide_json.get('img_keywords', '').strip()
        slide = presentation.slide_layouts[8]  # Picture with Caption
        slide = presentation.slides.add_slide(slide)
        placeholders = None

        title_placeholder = slide.shapes.title  # type: ignore[attr-defined]
        title_placeholder.text = self._remove_slide_number_from_heading(
            slide_json['heading']
        )

        try:
            pic_col = slide.shapes.placeholders[1]  # type: ignore[attr-defined]
        except KeyError:
            # Get placeholders from the slide without layout_number
            placeholders = self._get_slide_placeholders(slide)  # type: ignore[arg-type]
            pic_col = None
            for idx, name in placeholders:
                if 'picture' in name:
                    pic_col = slide.shapes.placeholders[idx]  # type: ignore[attr-defined]

        try:
            text_col = slide.shapes.placeholders[2]  # type: ignore[attr-defined]
        except KeyError:
            text_col = None
            if not placeholders:
                placeholders = self._get_slide_placeholders(slide)  # type: ignore[arg-type]

            for idx, name in placeholders:
                if 'content' in name:
                    text_col = slide.shapes.placeholders[idx]  # type: ignore[attr-defined]

        flat_items_list = self._get_flat_list_of_contents(
            slide_json['bullet_points'], level=0
        )
        self._add_bulleted_items(text_col.text_frame, flat_items_list)

        if not img_keywords:
            return True

        if isinstance(img_keywords, str) and img_keywords.startswith(
            ('http://', 'https://')
        ):
            try:
                img_response = requests.get(img_keywords, timeout=30)
                img_response.raise_for_status()
                image_data = BytesIO(img_response.content)
                pic_col.insert_picture(image_data)
                return True
            except Exception as ex:
                logger.error(
                    'Error while downloading image from URL: %s', str(ex)
                )

        try:
            url = 'https://api.pexels.com/v1/search'
            api_key = os.getenv('PEXELS_API_KEY')

            headers = {
                'Authorization': api_key,
                'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64; rv:10.0) '
                'Gecko/20100101 Firefox/10.0',
            }
            params = {
                'query': img_keywords,
                'size': 'medium',
                'page': 1,
                'per_page': 3,
            }
            response = requests.get(
                url, headers=headers, params=params, timeout=12
            )
            response.raise_for_status()
            json_response = response.json()

            if json_response.get('photos'):
                photo = random.choice(json_response['photos'])
                photo_url = photo.get('src', {}).get('large') or photo.get(
                    'src', {}
                ).get('original')

                if photo_url:
                    # Download and insert the image
                    img_response = requests.get(
                        photo_url, headers=headers, stream=True, timeout=12
                    )
                    img_response.raise_for_status()
                    image_data = BytesIO(img_response.content)

                    pic_col.insert_picture(image_data)
        except Exception as ex:
            logger.error(
                'Error occurred while adding image to slide: %s', str(ex)
            )

        return True

    def _handle_table(
        self,
        presentation: "presentation.Presentation",
        slide_json: Dict[str, Any],
    ) -> None:
        r"""Add a table to a slide.

        Args:
            presentation (presentation.Presentation): The presentation object.
            slide_json (Dict[str, Any]): The content of the slide as JSON data.
        """
        headers = slide_json['table'].get('headers', [])
        rows = slide_json['table'].get('rows', [])
        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        shapes.title.text = self._remove_slide_number_from_heading(
            slide_json['heading']
        )
        left = slide.placeholders[1].left
        top = slide.placeholders[1].top
        width = slide.placeholders[1].width
        height = slide.placeholders[1].height
        table = slide.shapes.add_table(
            len(rows) + 1, len(headers), left, top, width, height
        ).table

        # Set headers
        for col_idx, header_text in enumerate(headers):
            table.cell(0, col_idx).text = header_text
            table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True

        # Fill in rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_text in enumerate(row_data):
                table.cell(row_idx, col_idx).text = cell_text

    def _handle_step_by_step_process(
        self,
        presentation: "presentation.Presentation",
        slide_json: Dict[str, Any],
        slide_width_inch: float,
        slide_height_inch: float,
    ) -> None:
        r"""Add shapes to display a step-by-step process in the slide.

        Args:
            presentation (presentation.Presentation): The presentation object.
            slide_json (Dict[str, Any]): The content of the slide as JSON data.
            slide_width_inch (float): The width of the slide in inches.
            slide_height_inch (float): The height of the slide in inches.
        """
        import pptx
        from pptx.util import Inches

        steps = slide_json['bullet_points']
        n_steps = len(steps)

        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        shapes.title.text = self._remove_slide_number_from_heading(
            slide_json['heading']
        )

        if 3 <= n_steps <= 4:
            # Horizontal display
            height = Inches(1.5)
            width = Inches(slide_width_inch / n_steps - 0.01)
            top = Inches(slide_height_inch / 2)
            left = Inches(
                (slide_width_inch - width.inches * n_steps) / 2 + 0.05
            )

            for step in steps:
                shape = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height
                )
                text_frame = shape.text_frame
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = pptx.enum.text.PP_ALIGN.CENTER
                text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
                self._format_text(
                    paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                )
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                left = Inches(left.inches + width.inches - Inches(0.4).inches)
        elif 4 < n_steps <= 6:
            # Vertical display
            height = Inches(0.65)
            top = Inches(slide_height_inch / 4)
            left = Inches(1)
            width = Inches(slide_width_inch * 2 / 3)

            for step in steps:
                shape = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.PENTAGON, left, top, width, height
                )
                text_frame = shape.text_frame
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = pptx.enum.text.PP_ALIGN.CENTER
                text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
                self._format_text(
                    paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                )
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                top = Inches(top.inches + height.inches + Inches(0.3).inches)
                left = Inches(left.inches + Inches(0.5).inches)

    def _handle_icons_ideas(
        self,
        presentation: "presentation.Presentation",
        slide_json: Dict[str, Any],
        slide_width_inch: float,
        slide_height_inch: float,
    ) -> bool:
        r"""Add a slide with some icons and text.
        Uses semantic search to find the closest PNG icon from
        examples/toolkits/png128. If not found, falls back to first letter.
        """
        import os

        from pptx.dml.color import RGBColor
        from pptx.enum import dml
        from pptx.util import Inches

        if slide_json.get('bullet_points'):
            items = slide_json['bullet_points']
            for step in items:
                if not isinstance(step, str) or not step.startswith(
                    ICON_BEGINNING_MARKER
                ):
                    return False
            slide_layout = presentation.slide_layouts[5]
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = self._remove_slide_number_from_heading(
                slide_json['heading']
            )
            n_items = len(items)
            text_box_size = Inches(2)
            total_width = n_items * ICON_SIZE
            spacing = (Inches(slide_width_inch) - total_width) / (n_items + 1)
            top = Inches(3)
            icons_texts = []
            for item in items:
                match = ICONS_REGEX.search(item)
                if match is not None:
                    icons_texts.append((match.group(1), match.group(2)))
                else:
                    icons_texts.append(("", ""))
            for idx, item in enumerate(icons_texts):
                icon_query, accompanying_text = item
                left = spacing + idx * (ICON_SIZE + spacing)
                center = left + ICON_SIZE / 2
                shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    center - Inches(0.5),
                    top - (ICON_BG_SIZE - ICON_SIZE) / 2,
                    ICON_BG_SIZE,
                    ICON_BG_SIZE,
                )
                shape.fill.solid()
                shape.shadow.inherit = False
                # Use a random color for each icon background
                r = random.randint(0, 255)
                g = random.randint(0, 255)
                b = random.randint(0, 255)
                color = RGBColor(r, g, b)
                shape.fill.fore_color.rgb = shape.line.color.rgb = color
                # Semantic search for PNG icon
                icon_path = self._find_closest_png_icon(icon_query)
                inserted_icon = False
                if icon_path and os.path.exists(icon_path):
                    try:
                        slide.shapes.add_picture(
                            icon_path, left, top, height=ICON_SIZE
                        )
                        inserted_icon = True
                    except Exception as e:
                        logger.warning(
                            f"Failed to insert PNG icon '{icon_query}': {e!s}"
                        )
                        inserted_icon = False
                if not inserted_icon:
                    text_frame = shape.text_frame
                    text_frame.word_wrap = True
                    text_frame.paragraphs[
                        0
                    ].alignment = pptx.enum.text.PP_ALIGN.CENTER
                    text_frame.vertical_anchor = (
                        pptx.enum.text.MSO_ANCHOR.MIDDLE
                    )
                    run = text_frame.paragraphs[0].add_run()
                    run.text = icon_query[0].upper()
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                text_box = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    left=center - text_box_size / 2,
                    top=top + ICON_SIZE + Inches(0.2),
                    width=text_box_size,
                    height=text_box_size,
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.paragraphs[
                    0
                ].alignment = pptx.enum.text.PP_ALIGN.CENTER
                self._format_text(text_frame.paragraphs[0], accompanying_text)
                text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
                text_box.fill.background()
                text_box.line.fill.background()
                text_box.shadow.inherit = False
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.theme_color = dml.MSO_THEME_COLOR.TEXT_2
            return True
        return False

    def _remove_slide_number_from_heading(self, header: str) -> str:
        r"""Remove the slide number from a given slide header.

        Args:
            header (str): The header of a slide.

        Returns:
            str: The header without slide number.
        """
        if SLIDE_NUMBER_REGEX.match(header):
            idx = header.find(':')
            header = header[idx + 1 :]
        return header

    def _get_slide_placeholders(
        self,
        slide: "Slide",
    ) -> List[Tuple[int, str]]:
        r"""Return the index and name of all placeholders present in a slide.

        Args:
            slide (Slide): The slide.

        Returns:
            List[Tuple[int, str]]: A list containing placeholders (idx, name)
                tuples.
        """
        if hasattr(slide.shapes, 'placeholders'):
            placeholders = [
                (shape.placeholder_format.idx, shape.name.lower())
                for shape in slide.shapes.placeholders
            ]
            if placeholders and len(placeholders) > 0:
                placeholders.pop(0)  # Remove the title placeholder
            return placeholders
        return []

    def add_slides(
        self,
        file_path: str,
        content: str,
    ) -> str:
        r"""Add one or more slides to an existing PowerPoint presentation.

        Args:
            file_path (str): The path to the existing PPTX file.
            content (str): The content to add as a JSON string. Must represent
            a list of dictionaries with the following structure:
                - Content slides, which can be one of:
                    * Bullet/step slides: {
                        "heading": str,
                        "bullet_points": list of str or nested lists,
                        "img_keywords": str (optional)
                      }
                        - If any bullet point starts with '>> ', it will be
                          rendered as a step-by-step process.
                        - "img_keywords" can be a URL or search keywords for
                        an image (optional).
                    * Table slides: {
                        "heading": str,
                        "table": {
                            "headers": list of str,
                            "rows": list of list of str
                        }
                      }

        Returns:
            str: A success message indicating the slides were added.

        Example:
        [
            {
                "heading": "New Slide Title",
                "bullet_points": [
                    "**Bold text** for emphasis",
                    "*Italic text* for additional emphasis",
                    "Regular text for normal content"
                ],
                "img_keywords": "relevant search terms for images"
            },
            {
                "heading": "New Table Slide",
                "table": {
                    "headers": ["Column 1", "Column 2"],
                    "rows": [
                        ["Row 1, Col 1", "Row 1, Col 2"],
                        ["Row 2, Col 1", "Row 2, Col 2"]
                    ]
                }
            }
        ]
        """
        import pptx

        # Resolve file path
        resolved_path = self._resolve_filepath(file_path)

        if not resolved_path.exists():
            return f"File not found: {resolved_path}"

        # Parse and validate content format
        try:
            import json

            parsed_content = json.loads(content)
        except json.JSONDecodeError as e:
            logger.error(f"Content must be valid JSON: {e}")
            return "Failed to parse content as JSON"

        if not isinstance(parsed_content, list):
            logger.error(
                f"PPTX content must be a list of dictionaries, "
                f"got {type(parsed_content).__name__}"
            )
            return "PPTX content must be a list of dictionaries"

        try:
            # Load existing presentation
            presentation = pptx.Presentation(str(resolved_path))
            slide_width_inch, slide_height_inch = (
                self._get_slide_width_height_inches(presentation)
            )

            # Process new slides
            for slide_data in parsed_content:
                if not isinstance(slide_data, dict):
                    continue
                if 'table' in slide_data:
                    self._handle_table(presentation, slide_data)
                elif 'bullet_points' in slide_data:
                    if any(
                        step.startswith(STEP_BY_STEP_PROCESS_MARKER)
                        for step in slide_data['bullet_points']
                    ):
                        self._handle_step_by_step_process(
                            presentation,
                            slide_data,
                            slide_width_inch,
                            slide_height_inch,
                        )
                    else:
                        self._handle_default_display(presentation, slide_data)

            # Save the presentation
            presentation.save(str(resolved_path))
            success_msg = (
                f"Successfully added {len(parsed_content)} slides to "
                f"{resolved_path}"
            )
            logger.info(success_msg)
            return success_msg

        except Exception as e:
            error_msg = f"Failed to add slides to {resolved_path}: {e!s}"
            logger.error(error_msg)
            return error_msg

    def get_tools(self) -> List[FunctionTool]:
        r"""Returns a list of FunctionTool objects representing the
        functions in the toolkit.

        Returns:
            List[FunctionTool]: A list of FunctionTool objects
                representing the functions in the toolkit.
        """
        return [
            FunctionTool(self.create_presentation),
            FunctionTool(self.add_slides),
        ]
