# ========= Copyright 2023-2024 @ CAMEL-AI.org. All Rights Reserved. =========
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
# ========= Copyright 2023-2024 @ CAMEL-AI.org. All Rights Reserved. =========

import os
import tempfile

import pytest

from camel.toolkits.pptx_toolkit import PPTXToolkit


@pytest.fixture
def pptx_toolkit():
    r"""Fixture that creates a PPTXToolkit instance with a temporary
    directory."""
    with tempfile.TemporaryDirectory() as temp_dir:
        toolkit = PPTXToolkit(working_directory=temp_dir)
        yield toolkit


def test_pptx_toolkit_initialization(pptx_toolkit):
    r"""Test that PPTXToolkit initializes correctly."""
    assert isinstance(pptx_toolkit, PPTXToolkit)
    assert pptx_toolkit.working_directory.exists()


def test_get_tools(pptx_toolkit):
    r"""Test that get_tools returns the correct number of tools."""
    tools = pptx_toolkit.get_tools()
    assert len(tools) == 1
    assert tools[0].get_function_name() == "create_presentation"


def test_create_presentation_basic(pptx_toolkit):
    r"""Test creating a PPTX file with basic slide structure."""
    import json

    content = [
        {
            "title": "Test Presentation",
            "subtitle": "Created by PPTXToolkit",
        },
        {
            "heading": "First Slide",
            "bullet_points": ["This is the content of the first slide."],
        },
        {
            "heading": "Second Slide",
            "bullet_points": ["This is the content of the second slide."],
        },
    ]
    filename = "test_basic.pptx"

    result = pptx_toolkit.create_presentation(json.dumps(content), filename)

    # Check the result message
    assert "successfully created" in result

    # Check the file exists
    file_path = pptx_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    # Verify it's a valid PPTX file by checking we can open it
    try:
        import pptx

        presentation = pptx.Presentation(str(file_path))
        # Should have 3 slides (1 title + 2 content)
        assert len(presentation.slides) == 3
    except ImportError:
        # If python-pptx is not available, just check file exists
        pass


def test_create_presentation_with_images(pptx_toolkit):
    r"""Test creating a PPTX file with image URLs."""
    import json
    import os
    from unittest.mock import patch

    content = [
        {
            "title": "Presentation with Images",
            "subtitle": "Testing image integration",
        },
        {
            "heading": "Slide with Image",
            "bullet_points": ["This slide contains an image."],
            "img_keywords": "https://via.placeholder.com/300x200.png",
        },
    ]
    filename = "test_images.pptx"

    with patch.dict(os.environ, {"PEXELS_API_KEY": "test_key"}):
        result = pptx_toolkit.create_presentation(
            json.dumps(content), filename
        )

    # Check the result message
    assert "successfully created" in result

    # Check the file exists
    file_path = pptx_toolkit._resolve_filepath(filename)
    assert file_path.exists()


def test_create_presentation_title_only(pptx_toolkit):
    r"""Test creating a PPTX file with only title slide."""
    import json

    content = [
        {
            "title": "Title Only Presentation",
            "subtitle": "No content slides",
        }
    ]
    filename = "test_title_only.pptx"

    result = pptx_toolkit.create_presentation(json.dumps(content), filename)

    # Check the result message
    assert "successfully created" in result

    # Check the file exists
    file_path = pptx_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    # Verify it's a valid PPTX file with 1 slide
    try:
        import pptx

        presentation = pptx.Presentation(str(file_path))
        assert len(presentation.slides) == 1
    except ImportError:
        # If python-pptx is not available, just check file exists
        pass


def test_create_presentation_missing_fields(pptx_toolkit):
    r"""Test creating a PPTX file with slides missing optional fields."""
    import json

    content = [
        {
            "title": "Presentation",
            # Missing subtitle
        },
        {
            "heading": "Slide without bullet points",
            "bullet_points": [],
            # Missing img_keywords
        },
        {
            "heading": "Basic slide",
            "bullet_points": ["Slide with text but minimal structure"],
        },
    ]
    filename = "test_missing_fields.pptx"

    result = pptx_toolkit.create_presentation(json.dumps(content), filename)

    # Should handle missing fields gracefully
    assert "successfully created" in result

    # Check the file exists
    file_path = pptx_toolkit._resolve_filepath(filename)
    assert file_path.exists()


def test_create_presentation_invalid_image_url(pptx_toolkit):
    r"""Test creating a PPTX file with invalid image URL (should continue
    without failing)."""
    import json
    import os
    from unittest.mock import patch

    content = [
        {
            "title": "Test Invalid Image",
            "subtitle": "Should handle gracefully",
        },
        {
            "heading": "Slide with Bad Image",
            "bullet_points": ["This slide has an invalid image URL."],
            "img_keywords": (
                "https://invalid-url-that-does-not-exist.example.com/"
                "image.jpg"
            ),
        },
    ]
    filename = "test_invalid_image.pptx"

    with patch.dict(os.environ, {"PEXELS_API_KEY": "test_key"}):
        result = pptx_toolkit.create_presentation(
            json.dumps(content), filename
        )

    # Should still succeed despite invalid image
    assert "successfully created" in result

    # Check the file exists
    file_path = pptx_toolkit._resolve_filepath(filename)
    assert file_path.exists()


def test_create_presentation_invalid_content_type(pptx_toolkit):
    r"""Test that creating a PPTX file with invalid content type raises
    ValueError."""
    import json

    filename = "test_invalid.pptx"

    # Test with invalid JSON string
    result = pptx_toolkit.create_presentation("invalid json content", filename)
    assert result == "Failed to parse content as JSON"

    # Test with JSON that's not a list
    result = pptx_toolkit.create_presentation(json.dumps({}), filename)
    assert result == "PPTX content must be a list of dictionaries"


def test_create_presentation_auto_extension(pptx_toolkit):
    r"""Test that the .pptx extension is automatically added if missing."""
    import json

    content = [
        {
            "title": "Test Auto Extension",
            "subtitle": "Should add .pptx extension",
        }
    ]
    filename = "test_auto_extension"  # No .pptx extension

    result = pptx_toolkit.create_presentation(json.dumps(content), filename)

    # Check the result message
    assert "successfully created" in result

    # Check the file exists with .pptx extension
    file_path = pptx_toolkit._resolve_filepath(filename + '.pptx')
    assert file_path.exists()


def test_sanitize_and_resolve_filepath(pptx_toolkit):
    r"""Test that _resolve_filepath sanitizes filenames with spaces and
    special characters."""
    # Test filename with spaces and special characters
    filename = "test file with spaces & special chars!.pptx"
    resolved_path = pptx_toolkit._resolve_filepath(filename)

    # Should replace special characters with underscores
    assert "test_file_with_spaces_special_chars_.pptx" in str(resolved_path)

    # Should be under the output directory
    assert str(pptx_toolkit.working_directory) in str(resolved_path)


def test_resolve_filepath_absolute_path(pptx_toolkit):
    r"""Test _resolve_filepath with absolute paths."""
    # Create a temporary file path
    with tempfile.NamedTemporaryFile(
        suffix='.pptx', delete=False
    ) as temp_file:
        temp_path = temp_file.name

    try:
        resolved_path = pptx_toolkit._resolve_filepath(temp_path)

        # Should return the absolute path as-is (after sanitization)
        assert resolved_path.is_absolute()
        assert resolved_path.suffix == '.pptx'
    finally:
        # Clean up
        if os.path.exists(temp_path):
            os.unlink(temp_path)


def test_create_presentation_empty_content(pptx_toolkit):
    r"""Test creating a PPTX file with empty content list."""
    import json

    content = []
    filename = "test_empty.pptx"

    result = pptx_toolkit.create_presentation(json.dumps(content), filename)

    # Should handle empty content gracefully
    assert "successfully created" in result

    # Check the file exists
    file_path = pptx_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    # Verify it's a valid PPTX file with no slides
    try:
        import pptx

        presentation = pptx.Presentation(str(file_path))
        assert len(presentation.slides) == 0
    except ImportError:
        # If python-pptx is not available, just check file exists
        pass
