# ========= Copyright 2023-2024 @ CAMEL-AI.org. All Rights Reserved. =========
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
# ========= Copyright 2023-2024 @ CAMEL-AI.org. All Rights Reserved. =========

import json
import os
import tempfile
from pathlib import Path

import pytest
import yaml

from camel.toolkits import FileWriteToolkit


@pytest.fixture
def file_write_toolkit():
    r"""Create a FileWriteToolkit instance for testing."""
    with tempfile.TemporaryDirectory() as temp_dir:
        toolkit = FileWriteToolkit(output_dir=temp_dir)
        yield toolkit


@pytest.fixture
def temp_dir():
    r"""Create a temporary directory for file operations."""
    with tempfile.TemporaryDirectory() as temp_dir:
        yield temp_dir


def test_initialization():
    r"""Test the initialization of FileWriteToolkit with default parameters."""
    toolkit = FileWriteToolkit()

    assert toolkit.output_dir == Path("./").resolve()
    assert toolkit.default_encoding == "utf-8"
    assert toolkit.backup_enabled is True


def test_initialization_with_custom_parameters():
    r"""Test the initialization of FileWriteToolkit with custom parameters."""
    output_dir = "./custom_dir"
    default_encoding = "latin-1"
    backup_enabled = False

    toolkit = FileWriteToolkit(
        output_dir=output_dir,
        default_encoding=default_encoding,
        backup_enabled=backup_enabled,
    )

    assert toolkit.output_dir == Path(output_dir).resolve()
    assert toolkit.default_encoding == default_encoding
    assert toolkit.backup_enabled is backup_enabled


def test_write_to_file_text(file_write_toolkit):
    r"""Test writing plain text to a file."""
    content = "Hello, world!"
    filename = "test.txt"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = f.read()

    assert file_content == content


def test_write_to_file_markdown(file_write_toolkit):
    r"""Test writing markdown content to a file."""
    content = "# Heading\n\nThis is a paragraph with **bold** text."
    filename = "test.md"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = f.read()

    assert file_content == content


def test_write_to_file_json(file_write_toolkit):
    r"""Test writing JSON data to a file."""
    data = {"name": "John", "age": 30, "city": "New York"}
    filename = "test.json"

    result = file_write_toolkit.write_to_file(data, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = json.load(f)

    assert file_content == data


def test_write_to_file_yaml(file_write_toolkit):
    r"""Test writing YAML data to a file."""
    data = {"name": "John", "age": 30, "city": "New York"}
    filename = "test.yaml"

    result = file_write_toolkit.write_to_file(data, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = yaml.safe_load(f)

    assert file_content == data


def test_write_to_file_csv_string(file_write_toolkit):
    r"""Test writing CSV content as a string to a file."""
    content = "name,age,city\nJohn,30,New York\nJane,25,San Francisco"
    filename = "test.csv"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = f.read()

    assert file_content == content


def test_write_to_file_csv_list(file_write_toolkit):
    r"""Test writing CSV content as a list of lists to a file."""
    content = [
        ["name", "age", "city"],
        ["John", 30, "New York"],
        ["Jane", 25, "San Francisco"],
    ]
    filename = "test_list.csv"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    # Read the CSV file and check content
    with open(file_path, "r", encoding="utf-8") as f:
        lines = f.readlines()

    assert lines[0].strip() == "name,age,city"
    assert lines[1].strip() == "John,30,New York"
    assert lines[2].strip() == "Jane,25,San Francisco"


def test_write_to_file_html(file_write_toolkit):
    r"""Test writing HTML content to a file."""
    content = "<html><body><h1>Hello, world!</h1></body></html>"
    filename = "test.html"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = f.read()

    assert file_content == content


def test_write_to_file_no_extension(file_write_toolkit):
    r"""Test writing content to a file with no extension (should use default
    format).
    """
    content = "# Default format test"
    filename = "test_no_extension"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists with default extension (.md) and has correct
    # content
    file_path = file_write_toolkit._resolve_filepath(filename + ".md")
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = f.read()

    assert file_content == content


def test_write_to_file_custom_encoding(file_write_toolkit):
    r"""Test writing content with a custom encoding."""
    content = "Hello, world with custom encoding!"
    filename = "test_encoding.txt"
    encoding = "latin-1"

    result = file_write_toolkit.write_to_file(
        content, filename, encoding=encoding
    )

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content when read with the same
    # encoding
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding=encoding) as f:
        file_content = f.read()

    assert file_content == content


def test_write_to_file_nested_directory(file_write_toolkit):
    r"""Test writing to a file in a nested directory that doesn't exist yet."""
    content = "Content in nested directory"
    filename = "nested/dir/test.txt"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    with open(file_path, "r", encoding="utf-8") as f:
        file_content = f.read()

    assert file_content == content


def test_write_to_file_absolute_path(temp_dir):
    r"""Test writing to a file using an absolute path."""
    toolkit = FileWriteToolkit(output_dir="./default")
    content = "Content with absolute path"
    filename = os.path.join(temp_dir, "absolute_path.txt")

    result = toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists and has correct content
    assert os.path.exists(filename)

    with open(filename, "r", encoding="utf-8") as f:
        file_content = f.read()

    assert file_content == content


def test_get_tools(file_write_toolkit):
    r"""Test that get_tools returns the correct function tools."""
    tools = file_write_toolkit.get_tools()

    # Check that we have the expected number of tools
    assert len(tools) == 1

    # Check that the tool has the correct function name
    assert tools[0].get_function_name() == "write_to_file"


def test_write_to_file_pptx_basic(file_write_toolkit):
    r"""Test writing PPTX content to a file with basic slide structure."""
    content = [
        {
            "title": "Test Presentation",
            "subtitle": "Created by FileWriteToolkit",
        },
        {
            "title": "First Slide",
            "text": "This is the content of the first slide.",
        },
        {
            "title": "Second Slide", 
            "text": "This is the content of the second slide.",
        },
    ]
    filename = "test_basic.pptx"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    # Verify it's a valid PPTX file by checking we can open it
    try:
        import pptx
        presentation = pptx.Presentation(str(file_path))
        # Should have 3 slides (1 title + 2 content)
        assert len(presentation.slides) == 3
    except ImportError:
        # If python-pptx is not available, just check file exists
        pass


def test_write_to_file_pptx_with_images(file_write_toolkit):
    r"""Test writing PPTX content with image URLs."""
    content = [
        {
            "title": "Presentation with Images",
            "subtitle": "Testing image integration",
        },
        {
            "title": "Slide with Image",
            "text": "This slide contains an image.",
            "image": "https://via.placeholder.com/300x200.png",
        },
    ]
    filename = "test_images.pptx"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

def test_write_to_file_pptx_title_only(file_write_toolkit):
    r"""Test writing PPTX with only title slide."""
    content = [
        {
            "title": "Title Only Presentation",
            "subtitle": "No content slides",
        }
    ]
    filename = "test_title_only.pptx"

    result = file_write_toolkit.write_to_file(content, filename)

    # Check the result message
    assert "successfully written" in result

    # Check the file exists
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()

    # Verify it's a valid PPTX file with 1 slide
    try:
        import pptx
        presentation = pptx.Presentation(str(file_path))
        assert len(presentation.slides) == 1
    except ImportError:
        # If python-pptx is not available, just check file exists
        pass


def test_write_to_file_pptx_missing_fields(file_write_toolkit):
    r"""Test writing PPTX with slides missing optional fields."""
    content = [
        {
            "title": "Presentation",
            # Missing subtitle
        },
        {
            "title": "Slide without text",
            # Missing text and image
        },
        {
            # Missing title
            "text": "Slide with text but no title",
        },
    ]
    filename = "test_missing_fields.pptx"

    result = file_write_toolkit.write_to_file(content, filename)

    # Should handle missing fields gracefully
    assert "successfully written" in result

    # Check the file exists
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()


def test_write_to_file_pptx_invalid_image_url(file_write_toolkit):
    r"""Test writing PPTX with invalid image URL (should continue without failing)."""
    content = [
        {
            "title": "Test Invalid Image",
            "subtitle": "Should handle gracefully",
        },
        {
            "title": "Slide with Bad Image",
            "text": "This slide has an invalid image URL.",
            "image": "https://invalid-url-that-does-not-exist.example.com/image.jpg",
        },
    ]
    filename = "test_invalid_image.pptx"

    result = file_write_toolkit.write_to_file(content, filename)

    # Should still succeed despite invalid image
    assert "successfully written" in result

    # Check the file exists
    file_path = file_write_toolkit._resolve_filepath(filename)
    assert file_path.exists()


def test_sanitize_and_resolve_filepath(file_write_toolkit):
    r"""Test that _resolve_filepath sanitizes filenames
    with spaces and special characters.
    """
    # Filename with spaces and special characters
    unsafe_filename = "My Video: How to Fix! File @ Name?.md"
    # Expected sanitized filename (all disallowed characters replaced by
    # underscores)
    expected_sanitized = file_write_toolkit._sanitize_filename(unsafe_filename)

    # Resolve the filepath using the toolkit
    resolved_path = file_write_toolkit._resolve_filepath(unsafe_filename)
    # Expected path is in the toolkit's output_dir with the sanitized filename
    expected_path = file_write_toolkit.output_dir / expected_sanitized

    # Check that the resolved path matches the expected path
    assert (
        resolved_path == expected_path.resolve()
    ), "The resolved file path does not match the expected sanitized path."
